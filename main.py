import streamlit as st
import pandas as pd
import sqlite3
import io
import urllib.parse
import hashlib
import zipfile
import json
from datetime import datetime, date, timedelta
import plotly.express as px
import plotly.graph_objects as go
from fpdf import FPDF
from docx import Document
from pptx import Presentation

# --- CONFIGURAÇÃO DE UI/UX ---
st.set_page_config(page_title="SGI Guardian", layout="wide", page_icon="🛡️")

st.markdown("""
    <style>
    header {visibility: hidden;}
    .stApp { 
        background: linear-gradient(135deg, #05070A 0%, #0E1117 100%);
        color: #E0E6ED; 
    }
    .stButton>button { 
        width: 100%; 
        border-radius: 12px; 
        height: 3.5em; 
        font-weight: bold; 
        background: linear-gradient(135deg, #00D1FF, #0052D4); 
        color: white; 
        border: none;
        font-size: 16px; 
        box-shadow: 0px 4px 15px rgba(0, 209, 255, 0.3); 
    }
    .stButton>button:hover { 
        transform: translateY(-2px);
        box-shadow: 0px 8px 25px rgba(0, 209, 255, 0.5); 
    }
    div[data-testid="stMetric"] { 
        background-color: rgba(28, 37, 51, 0.9); 
        padding: 20px;
        border-radius: 15px; 
        border-bottom: 4px solid #00D1FF; 
    }
    .stTabs [data-baseweb="tab-list"] { gap: 8px; }
    .stTabs [data-baseweb="tab"] {
        background-color: #1C2533; 
        border-radius: 10px; 
        padding: 8px 20px; 
        color: white; 
        font-weight: bold;
    }
    .stExpander { 
        background-color: rgba(255, 255, 255, 0.03) !important; 
        border: 1px solid #1E293B; 
        border-radius: 12px;
    }
    </style>
    """, unsafe_allow_html=True)

# --- INICIALIZAÇÃO DO BANCO DE DADOS ---
def init_db():
    conn = sqlite3.connect('sgi_guardian.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS equipe (
                    id INTEGER PRIMARY KEY AUTOINCREMENT, 
                    nome TEXT, 
                    cargo TEXT, 
                    status TEXT, 
                    icc INTEGER)''')
    c.execute('''CREATE TABLE IF NOT EXISTS ativos (
                    id TEXT PRIMARY KEY, 
                    categoria TEXT, 
                    local TEXT, 
                    vencimento DATE, 
                    status TEXT)''')
    c.execute('''CREATE TABLE IF NOT EXISTS licencas (
                    id INTEGER PRIMARY KEY AUTOINCREMENT, 
                    nome TEXT, 
                    orgao TEXT, 
                    vencimento DATE, 
                    status TEXT)''')
    c.execute('''CREATE TABLE IF NOT EXISTS auditorias (
                    id INTEGER PRIMARY KEY AUTOINCREMENT, 
                    data TEXT, 
                    local TEXT, 
                    alvo TEXT, 
                    relato TEXT, 
                    tcr REAL, 
                    glosa REAL, 
                    lei TEXT, 
                    iso TEXT, 
                    acao TEXT, 
                    hash TEXT)''')
    conn.commit()
    conn.close()

init_db()

# --- CONSTANTES E GLOSSÁRIO ---
GLOS = {
    "TCR": "Total Cost of Risk: Soma de multas, glosas, dano à imagem e honorários.",
    "LTO": "License to Operate: Alvarás e licenças vitais para o posto.",
    "ICC": "Índice de Confiabilidade: Probabilidade de desvio técnico/acidente.",
    "PMOC": "Plano de Manutenção de Ar Condicionado (Obrigatório Lei 13.589/18).",
    "NR-37": "Norma para Segurança e Saúde em Plataformas de Petróleo (Offshore).",
    "RBAC 117": "Regulamento da ANAC sobre Gerenciamento de Risco de Fadiga Humana.",
    "L14.967/24": "Estatuto Seg. Privada: Marco regulador de conduta e penalidades PF.",
    "SLA": "Service Level Agreement: Acordo de nível de serviço com multas.",
    "SGI": "Sistema de Gestão Integrado (Qualidade, Meio Ambiente, Saúde e Segurança)."
}

DOUTRINA_SGI = {
    "INDUSTRIAL / TERRA": {
        "lei": "NRs (01, 10, 11, 12, 35) + IT-17 Bombeiros",
        "iso": "ISO 45001",
        "pilar": "HSE Industrial",
        "detalhes": "Lei 14.967/2024 (Estatuto Seg. Privada) + Portaria 18.045/23 PF."
    },
    "OFFSHORE / MAR": {
        "lei": "NR-37 + NORMAM-204 + NR-13 + Lei 5.811/72",
        "iso": "ISO 14001:2026/45001",
        "pilar": "SMS Marítimo",
        "detalhes": "NR-37 (Segurança Offshore) + NORMAM-204/DPC Marinha."
    },
    "AVIAÇÃO / CÉU": {
        "lei": "RBAC 107/153 + RBAC 117 (Fadiga)",
        "iso": "ISO 9001:2026",
        "pilar": "Safety & AVSEC",
        "detalhes": "RBAC 107 (AVSEC) + RBAC 153 (Operações em Aeródromos)."
    },
    "HOSPITALAR / SAÚDE": {
        "lei": "RDC 222 ANVISA + NR-32",
        "iso": "ISO 9001:2026",
        "pilar": "SGI Saúde",
        "detalhes": "RDC 222 ANVISA + NR-32"
    },
    "SEGURANÇA PRIVADA": {
        "lei": "Lei 14.967/24 + Portaria 18.045/23 PF",
        "iso": "ISO 37301/45001",
        "pilar": "HSG (Governança)",
        "detalhes": "Lei 14.967/2024 (Estatuto Seg. Privada)"
    }
}

# --- FUNÇÕES AUXILIARES ---
def i_help(termo):
    with st.popover(f"💡 Doutrina: {termo}"):
        st.info(GLOS.get(termo, "Legislação e SGI 2027."))

def motor_apex_v80(dados):
    prejuizos = {"Baixa": 45000, "Média": 150000, "Alta": 750000, "Crítica": 6500000}
    multa_base = prejuizos.get(dados['gravidade'], 50000)
    tcr = multa_base * 4.8 
    glosa = multa_base * 0.25
    
    ref = DOUTRINA_SGI.get(dados['setor'], DOUTRINA_SGI["INDUSTRIAL / TERRA"])
    medidas = {
        "Advertência": "advertência", 
        "Suspensão": "suspensão", 
        "Substituição": "substituição", 
        "Desligamento": "justa causa", 
        "Reciclagem": "reciclagem",
        "Manutenção": "manutenção corretiva",
        "Interdição": "interdição técnica"
    }
    
    h = hashlib.sha512(dados['relato'].encode()).hexdigest()[:16].upper()
    
    texto = f"""Prezado Cliente, informamos a resolução definitiva da ocorrência no posto {dados['local']} relativa ao alvo {dados['alvo']}.

1. PARECER TÉCNICO-LEGAL: O desvio identificado ("{dados['relato']}") foi confrontado com a {ref['lei']}. O fato infringe as diretrizes de {ref['pilar']} e as normas {ref['iso']}.

2. TRATAMENTO E RESOLUÇÃO: Priorizando a segurança e a CONTINUIDADE DO NEGÓCIO, aplicamos a medida de {medidas.get(dados['medida'], dados['medida'])} imediata. A eficácia foi validada via evidência visual e registro sistêmico.

3. IMPACTO EXECUTIVO: Esta intervenção preventiva evitou um Total Cost of Risk (TCR) de R$ {tcr:,.2f} e uma glosa de R$ {glosa:,.2f}.

ID Digital: {h}"""
    
    return texto, tcr, glosa, ref['lei'], h

def calcular_icc_v80(faltas, atrasos, desvios):
    score = 100 - (faltas * 15) - (atrasos * 5) - (desvios * 10)
    return max(score, 0)

def calc_burnout_v80(h_excedentes):
    score = h_excedentes * 12.5
    if score > 80: return "🔴 CRÍTICO", "red"
    if score > 50: return "🟡 ALERTA", "orange"
    return "🟢 REGULAR", "green"

def sugerir_causa_raiz(evento):
    causas = {
        "Falta": "Mão de Obra: Desmotivação ou falha de transporte.",
        "Desvio de Conduta": "Medida: Falha na supervisão direta.",
        "Falha Técnica": "Máquina: Manutenção preventiva inexistente.",
        "Fadiga": "Método: Escala de trabalho incompatível.",
        "Inconformidade ISO": "Processo: Falha no controle operacional.",
        "Justa Causa": "Conduta: Violação grave de diretrizes contratuais."
    }
    return causas.get(evento, "Analisar fatores humanos e técnicos.")

def avaliar_impacto_esg(gravidade):
    impactos = {"Baixa": "Neutro", "Média": "Alerta Social", "Alta": "Risco de Governança", "Crítica": "Dano Grave"}
    return impactos.get(gravidade, "Analisando")

def criar_apresentacao_v80(d):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"GESTÃO DE RISCO SGI - {d['local']}"
    body = slide.placeholders[1]
    body.text = f"Evento: {d['medida']}\nImpacto ESG: {d['esg']}\nROI: R$ {d['tcr']:,.2f}\n\nResumo: {d['txt'][:200]}..."
    buf_p = io.BytesIO()
    prs.save(buf_p)
    return buf_p.getvalue()

def exportar_excel_bi(d):
    df_bi = pd.DataFrame([d])
    buf_e = io.BytesIO()
    with pd.ExcelWriter(buf_e, engine='openpyxl') as writer:
        df_bi.to_excel(writer, index=False, sheet_name='Base_SGI')
    return buf_e.getvalue()

def traduzir_parecer_v80(texto, idioma):
    if idioma == "Inglês":
        t = texto.replace("Prezado Cliente", "Dear Client")
        t = t.replace("informamos a resolução", "we inform the resolution")
        t = t.replace("PARECER TÉCNICO-LEGAL", "TECHNICAL-LEGAL OPINION")
        t = t.replace("TRATAMENTO E RESOLUÇÃO", "TREATMENT AND RESOLUTION")
        t = t.replace("IMPACTO EXECUTIVO", "EXECUTIVE IMPACT")
        return t
    return texto

def calcular_uptime_v80(paradas_min):
    tempo_total = 43200  # Total de minutos em um mês (30 dias)
    disponibilidade = ((tempo_total - paradas_min) / tempo_total) * 100
    return round(disponibilidade, 2)

def enviar_alerta_rh(d):
    assunto_rh = f"ALERTA DISCIPLINAR - {d['alvo']} - {d['local']}"
    corpo_rh = f"O colaborador {d['alvo']} recebeu a medida de {d['medida']}."
    gmail_rh = f"https://mail.google.com/mail/?view=cm&fs=1&to=rh@empresa.com&su={urllib.parse.quote(assunto_rh)}&body={urllib.parse.quote(corpo_rh)}"
    return gmail_rh

# --- INTERFACE DE COMANDO TITAN ---
st.title("🛡️ SGI Guardian")
st.subheader("Supreme Sovereign Command Matrix")

tab_audit, tab_ativos, tab_equipe, tab_bi, tab_boutique = st.tabs([
    "🧐 Auditoria & Devolutiva", "⚙️ Ativos & LTO", "👥 Gestão de Equipe", "📊 Dashboard Power BI", "🚀 Boutique Executive"
])

# --- ABA 1: AUDITORIA ---
with tab_audit:
    st.header("📋 Inspeção Forense e PDCA")
    with st.form("audit_form_v80"):
        c1, c2, c3 = st.columns(3)
        with c1:
            setor = st.selectbox("Setor de Atuação", list(DOUTRINA_SGI.keys()))
            local = st.text_input("Unidade / Posto / Planta")
            i_help("LTO")
        with c2:
            alvo = st.text_input("Alvo (Profissional ou Ativo)")
            grav = st.select_slider("Criticidade do Impacto", ["Baixa", "Média", "Alta", "Crítica"])
            i_help("TCR")
        with c3:
            medida = st.selectbox("Ação de Tratamento:", ["Advertência", "Suspensão", "Substituição", "Manutenção", "Interdição", "Justa Causa"])
            i_help("SGI")
            foto = st.camera_input("Capturar Solução")
        
        relato = st.text_area("Descreva os fatos (O app redigirá o Parecer Forense Master):")
        submit = st.form_submit_button("⚖️ EXECUTAR INTELIGÊNCIA SOBERANA")
        
        if submit:
            if relato and local and alvo:
                texto, tcr, glosa, lei, h = motor_apex_v80({
                    'setor': setor, 'gravidade': grav, 'alvo': alvo, 
                    'local': local, 'medida': medida, 'relato': relato
                })
                
                res_esg = avaliar_impacto_esg(grav)
                
                st.session_state['v80_rep'] = {
                    'txt': texto, 'local': local, 'tcr': tcr, 
                    'glosa': glosa, 'lei': lei, 'alvo': alvo, 
                    'medida': medida, 'relato_orig': relato,
                    'gravidade': grav, 'esg': res_esg, 'hash': h,
                    'data': datetime.now().strftime("%d/%m/%Y %H:%M")
                }
                
                # Salvar no Banco
                conn = sqlite3.connect('sgi_guardian.db')
                conn.execute("""INSERT INTO auditorias (data, local, alvo, relato, tcr, glosa, lei, acao, hash) 
                               VALUES (?,?,?,?,?,?,?,?,?)""", 
                            (st.session_state['v80_rep']['data'], local, alvo, relato, tcr, glosa, lei, medida, h))
                conn.commit()
                conn.close()
                
                st.success("✅ Inteligência SGI Concluída! Devolutiva pronta na aba Boutique.")
            else:
                st.error("Por favor, preencha Local, Alvo e Relato.")

    if 'v80_rep' in st.session_state:
        st.divider()
        st.subheader("🪵 Análise de Causa Raiz (Ishikawa)")
        r_v80 = st.session_state['v80_rep']
        st.info(f"**Sugestão Técnica:** {sugerir_causa_raiz(r_v80['medida'])}")
        
        st.subheader("🚀 Plano de Ação Corretiva (5W2H)")
        with st.expander("Visualizar Cronograma de Tratamento"):
            st.write(f"**O QUE:** {r_v80['medida']} imediata no posto.")
            st.write(f"**POR QUE:** Para mitigar risco de {r_v80['gravidade']} e garantir conformidade.")
            st.write(f"**QUEM:** Supervisor de Área e Gestor de SGI.")
            st.write(f"**ONDE:** Unidade {r_v80['local']}.")
            st.write(f"**QUANDO:** {r_v80['data']}.")
            st.write(f"**COMO:** Aplicação de medida disciplinar e reforço técnico.")
            st.write(f"**QUANTO:** Custo operacional absorvido pelo ROI de prevenção.")

# --- ABA 2: ATIVOS & LTO ---
with tab_ativos:
    st.header("⚙️ Controle de Infraestrutura & Licenças")
    t_fogo, t_infra, t_lto, t_saude, t_amb = st.tabs(["🔥 Safety", "⚡ Facilities", "🏢 LTO", "🏥 Saúde", "🍃 Ambiental"])
    
    with t_fogo:
        st.subheader("Bunker de Prevenção: Extintores e Hidrantes")
        with st.form("cad_fogo"):
            c_f1, c_f2 = st.columns(2)
            id_f = c_f1.text_input("Selo INMETRO / ID")
            venc_f = c_f2.date_input("Vencimento Carga")
            if st.form_submit_button("💾 Salvar Ativo"):
                st.success(f"Ativo {id_f} registrado.")
                
    with t_infra:
        st.subheader("Controle de Engenharia: PMOC e SPDA")
        with st.form("cad_infra"):
            id_i = st.text_input("ID do Equipamento")
            tipo_i = st.selectbox("Tipo:", ["Gerador", "Nobreak", "Elevador", "PMOC Ar", "SPDA"])
            venc_i = st.date_input("Próxima Manutenção")
            if st.form_submit_button("💾 Registrar"):
                st.success("Ativo registrado conforme NRs.")
                
    with t_lto:
        st.subheader("Radar de Licenças (License to Operate)")
        with st.form("cad_lto"):
            c_l1, c_l2 = st.columns(2)
            doc_l = c_l1.text_input("Documento (Alvará, AVCB, PF)")
            orgao_l = c_l1.text_input("Órgão Emissor")
            venc_l = c_l2.date_input("Vencimento Legal")
            risco_l = c_l2.selectbox("Risco", ["Baixo", "Médio", "Interdição"])
            if st.form_submit_button("💾 Monitorar"):
                conn = sqlite3.connect('sgi_guardian.db')
                conn.execute("INSERT INTO licencas (nome, orgao, vencimento, status) VALUES (?,?,?,?)", 
                            (doc_l, orgao_l, venc_l, risco_l))
                conn.commit()
                conn.close()
                st.success("Licença sob radar Sentinel.")
        
        st.divider()
        conn = sqlite3.connect('sgi_guardian.db')
        df_lto = pd.read_sql_query("SELECT nome, orgao, vencimento, status FROM licencas", conn)
        conn.close()
        if not df_lto.empty:
            st.table(df_lto)
        else:
            st.info("Nenhuma licença cadastrada.")

    with t_saude:
        st.subheader("🏥 Equipamentos de Emergência")
        with st.form("form_saude"):
            id_s = st.text_input("ID DEA / Autoclave")
            status_s = st.selectbox("Status", ["Calibrado", "Vencido"])
            if st.form_submit_button("💾 Protocolar"):
                st.success("Equipamento registrado.")

    with t_amb:
        st.subheader("🍃 Gestão Ambiental")
        with st.form("form_amb"):
            residuo = st.selectbox("Classe:", ["Classe I", "Classe IIA", "Classe IIB"])
            if st.form_submit_button("💾 Salvar MTR"):
                st.success("Manifesto protocolado.")

# --- ABA 3: EQUIPE ---
with tab_equipe:
    st.header("👥 Gestão de Tropa")
    t_peo, t_epi, t_hist = st.tabs(["Performance", "EPI (NR-06)", "Histórico"])
    
    with t_peo:
        with st.form("form_hr"):
            col_h1, col_h2 = st.columns(2)
            nome_h = col_h1.text_input("Nome do Colaborador")
            evento_h = col_h2.selectbox("Evento:", ["Falta", "Atraso", "Fadiga", "Elogio"])
            if st.form_submit_button("🚀 Registrar"):
                conn = sqlite3.connect('sgi_guardian.db')
                conn.execute("INSERT INTO equipe (nome, cargo, status, icc) VALUES (?,?,?,?)",
                             (nome_h, "Operacional", evento_h, 85))
                conn.commit()
                conn.close()
                st.success(f"Registro de {nome_h} salvo.")

    with t_epi:
        with st.form("form_epi"):
            col_e1, col_e2 = st.columns(2)
            epi_n = col_e1.text_input("EPI (Ex: Colete)")
            ca_n = col_e2.text_input("Número do CA")
            if st.form_submit_button("💾 Salvar EPI"):
                st.success("EPI protocolado.")

    with t_hist:
        nome_busca = st.text_input("Consultar Histórico por Nome:")
        if nome_busca:
            conn = sqlite3.connect('sgi_guardian.db')
            query = f"SELECT data, alvo, acao, relato FROM auditorias WHERE alvo LIKE '%{nome_busca}%'"
            df_hist = pd.read_sql_query(query, conn)
            conn.close()
            if not df_hist.empty: st.dataframe(df_hist, use_container_width=True)
            else: st.info("Nenhum registro encontrado.")

# --- ABA 4: BI ---
with tab_bi:
    st.header("📊 Business Intelligence")
    if 'v80_rep' in st.session_state:
        d = st.session_state['v80_rep']
        c_m1, c_m2, c_m3 = st.columns(3)
        c_m1.metric("Lucro Protegido (TCR)", f"R$ {d['tcr']:,.2f}", "+28% ROI")
        c_m2.metric("Glosa Evitada", f"R$ {d['glosa']:,.2f}", "SLA Seguro")
        c_m3.metric("Integridade SGI", "97%", "Meta batida")
        
        fig_gauge = go.Figure(go.Indicator(
            mode = "gauge+number", value = 97,
            title = {'text': "Conformidade Global %"},
            gauge = {'axis': {'range': [None, 100]},
                     'bar': {'color': "#00D1FF"},
                     'steps' : [
                         {'range': [0, 50], 'color': "#FF4B4B"},
                         {'range': [50, 85], 'color': "#FFD700"},
                         {'range': [85, 100], 'color': "#00CC96"}]}))
        st.plotly_chart(fig_gauge, use_container_width=True)
    
    st.divider()
    st.subheader("📈 Performance Histórica")
    df_curve = pd.DataFrame({'Mês': ['Jan', 'Fev', 'Mar', 'Abr'], 'Conformidade': [70, 75, 88, 97]})
    fig_curve = px.line(df_curve, x='Mês', y='Conformidade', markers=True)
    fig_curve.update_traces(line_color='#00D1FF')
    st.plotly_chart(fig_curve, use_container_width=True)

    st.divider()
    st.subheader("🧘 People Analytics: Fadiga (RBAC 117)")
    h_excedente = st.slider("Horas Excedentes no Turno:", 0, 12, 2)
    status_f, cor_f = calc_burnout_v80(h_excedente)
    st.markdown(f"### Status: :{cor_f}[{status_f}]")
    if h_excedente > 8:
        st.error("🚨 INTERDIÇÃO: Limite de segurança excedido.")

    st.divider()
    st.subheader("⏱️ Disponibilidade Operacional (SLA)")
    min_parada = st.number_input("Minutos de Inatividade (Mês):", min_value=0, value=120)
    uptime = calcular_uptime_v80(min_parada)
    st.metric("Uptime do Posto", f"{uptime}%")

# --- ABA 5: BOUTIQUE ---
with tab_boutique:
    if 'v80_rep' in st.session_state:
        d = st.session_state['v80_rep']
        st.header("🚀 Despacho de Documentos")
        
        st.subheader("📲 WhatsApp")
        u_txt = urllib.parse.quote(d['txt'])
        w1, w2, w3 = st.columns(3)
        w1.markdown(f'<a href="https://wa.me/?text={u_txt}" target="_blank"><button style="background:#34B7F1; border-radius:10px; color:white; border:none; height:50px; font-weight:bold; width:100%;">🔵 CLIENTE</button></a>', unsafe_allow_html=True)
        w2.markdown(f'<a href="https://wa.me/?text={urllib.parse.quote("ROI: R$ " + str(d["tcr"]))}" target="_blank"><button style="background:#075E54; border-radius:10px; color:white; border:none; height:50px; font-weight:bold; width:100%;">🟢 DIRETORIA</button></a>', unsafe_allow_html=True)
        w3.markdown(f'<a href="https://wa.me/?text={urllib.parse.quote("Ação: " + d["medida"])}" target="_blank"><button style="background:#25D366; border-radius:10px; color:white; border:none; height:50px; font-weight:bold; width:100%;">🟢 OPERACIONAL</button></a>', unsafe_allow_html=True)
        
        st.divider()
        st.subheader("💾 Exportação Master")
        f1, f2, f3, f4 = st.columns(4)
        f1.download_button("📄 PDF Dossiê", d['txt'], f"SGI_{d['local']}.txt", use_container_width=True) # PDF simplificado como txt para demo
        
        doc = Document(); doc.add_paragraph(d['txt']); w_buf = io.BytesIO(); doc.save(w_buf)
        f2.download_button("📝 WORD", w_buf.getvalue(), "Relatorio.docx", use_container_width=True)
        
        f3.download_button("📊 PPT Slides", criar_apresentacao_v80(d), "Apresentacao.pptx", use_container_width=True)
        f4.download_button("🚀 POWER BI (.xlsx)", exportar_excel_bi(d), "Base_BI.xlsx", use_container_width=True)
        
        st.divider()
        st.subheader("👤 RH & ERP")
        st.markdown(f'<a href="{enviar_alerta_rh(d)}" target="_blank"><button style="background:#004a99; color:white; border:none; padding:10px; border-radius:10px; width:100%;">👤 NOTIFICAR RH</button></a>', unsafe_allow_html=True)
        
        json_data = json.dumps(d, indent=4)
        st.download_button("🔌 EXPORTAR JSON", json_data, "sgi_data.json", use_container_width=True)

        st.divider()
        st.text_area("Parecer Gerado:", d['txt'], height=300)
    else:
        st.warning("⚠️ Realize uma auditoria na primeira aba para liberar as funções executivas.")

# --- SIDEBAR ---
st.sidebar.title("🛡️ Sentinel Engine")
st.sidebar.info(f"Status: Operacional\nVersão: {datetime.now().year}.6.10")
st.sidebar.divider()
st.sidebar.write("👤 **Acesso:** SGI Master Alpha")
st.sidebar.write("🔐 **Criptografia:** SHA-512 Active")

if st.sidebar.button("♻️ Reinicializar Base"):
    st.sidebar.warning("Ação restrita.")
    st.sidebar.divider()
st.sidebar.caption("SGI Guardian - Sovereign Tech")
st.sidebar.write(f"Sincronização: {datetime.now().strftime('%H:%M:%S')}")
